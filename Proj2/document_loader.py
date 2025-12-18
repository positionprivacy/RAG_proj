import os
import io
import uuid # 用于生成唯一临时文件名
import concurrent.futures # 用于多线程并发
from typing import List, Dict, Optional
import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI

# 引入多模态需要的库
import fitz  # PyMuPDF, 用于PDF图片提取
import dashscope
from http import HTTPStatus
import config  # 导入配置文件

class DocumentLoader:
    def __init__(
        self,
        data_dir: str = config.DATA_DIR,
    ):
        self.data_dir = data_dir
        # [修改] 增加图片格式支持
        self.code_map = {
            '.py': 'python', '.c': 'c', '.cpp': 'cpp', '.java': 'java', 
            '.js': 'javascript', '.ts': 'typescript', '.html': 'html', 
            '.css': 'css', '.sql': 'sql', '.go': 'go', '.rs': 'rust', 
            '.sh': 'bash', '.md': 'markdown', '.json': 'json', 
            '.xml': 'xml', '.yaml': 'yaml'
        }
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"] + \
                                 [".png", ".jpg", ".jpeg"] + \
                                 list(self.code_map.keys())
        
        # 配置 Dashscope API
        dashscope.api_key = config.OPENAI_API_KEY

        self.client = OpenAI(api_key=config.OPENAI_API_KEY, base_url=config.OPENAI_API_BASE)
        
        # 并发控制：设置最大工作线程数
        self.max_workers = 10

    def _describe_image(self, image_bytes: bytes, source_info: str, context_text: str = "") -> str:
        """
        [内部辅助函数] 调用 Qwen-VL 对图片进行描述
        包含并发安全(UUID)和上下文感知逻辑
        """
        if not hasattr(config, "VL_MODEL_NAME") or not config.VL_MODEL_NAME:
            return ""

        # 使用 uuid 生成唯一文件名，防止多线程冲突
        unique_name = f"temp_img_{uuid.uuid4().hex}.png"

        try:
            with open(unique_name, "wb") as f:
                f.write(image_bytes)
            
            # 只有大于 5KB 的图片才处理
            if os.path.getsize(unique_name) < 5 * 1024:
                return ""

            # print(f"  > [并发] 调用 Qwen-VL: {source_info}...") 
            
            # --- 构建包含上下文的 Prompt ---
            safe_context = context_text[:1000].replace("\n", " ") if context_text else "无"
            
            prompt_content = (
                f"这张图出现在课程课件或用户上传中。\n"
                f"【背景文字参考】：{safe_context}\n\n"
                f"请结合上下文描述图片内容。如果图片包含具体知识点（如架构图、公式、代码截图），请详细提取文字和含义；"
                f"如果图片只是装饰或与上下文无关，请简要概括或忽略。"
            )

            messages = [
                {
                    "role": "user",
                    "content": [
                        {"image": f"file://{os.path.abspath(unique_name)}"},
                        {"text": prompt_content}
                    ]
                }
            ]
            
            response = dashscope.MultiModalConversation.call(
                model=config.VL_MODEL_NAME,
                messages=messages
            )
            
            if response.status_code == HTTPStatus.OK:
                desc = response.output.choices[0].message.content[0]['text']
                return f"\n[图片内容描述]: {desc}\n"
            else:
                print(f"  ! 图片处理失败 ({source_info}): {response.message}")
                return ""
                
        except Exception as e:
            print(f"  ! 图片处理异常 ({source_info}): {e}")
            return ""
        finally:
            if os.path.exists(unique_name):
                os.remove(unique_name)

    def _generate_summary(self, text: str) -> str:
        """
        [内部辅助函数] 使用 LLM 生成文本摘要
        """
        if not text.strip():
            return ""
        
        system_instruction = "你是一位专业的总结助手。请将下面的完整文件内容总结为一个简洁、准确的段落，保留核心要点和技术术语，不要增加文件内容中不存在的内容。总结内容前面加上文件涉及的课程名称/三级学科名称，总结需概括、精炼，严格保证50字以内。"
        
        text_input = text[:6000] # 截断

        prompt = system_instruction + "\n\n请基于以下内容生成摘要：\n" + text_input.strip()

        messages = [
            {"role": "user", "content": prompt} 
        ]

        try:
            response = self.client.chat.completions.create(
                model=config.MODEL_NAME, messages=messages, temperature=0.5, max_tokens=100
            )
            summary = "片段摘要：" + response.choices[0].message.content
        except Exception as e:
            print(f" ! 摘要生成失败: {e}")
            return text[:100] + ("... (摘要失败)" if len(text) > 200 else "")
        
        return summary

    # ================= 核心加载逻辑 =================

    def load_image(self, file_path: str) -> str:
        """
        [新增] 专门处理纯图片文件 (.png, .jpg)
        """
        print(f"正在加载独立图片: {file_path}")
        try:
            with open(file_path, "rb") as f:
                image_bytes = f.read()
            
            # 独立图片没有上下文文本，传空字符串
            desc = self._describe_image(
                image_bytes, 
                source_info=os.path.basename(file_path),
                context_text="" 
            )
            return desc if desc else "[图片无法识别或内容为空]"
        except Exception as e:
            return f"[图片解析错误: {str(e)}]"

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，支持多模态并发加速"""
        print(f"正在加载 PDF (并发模式): {file_path}")
        results = []
        reader = PdfReader(file_path)
        
        doc_fitz = None
        try: doc_fitz = fitz.open(file_path)
        except: print("提示: 未安装 PyMuPDF，跳过图片识别")

        page_tasks = [] # (page_text, list_of_futures)
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            for i, page in enumerate(doc_fitz):
                page_text = page.get_text() or ""
                image_futures = []

                if doc_fitz and i < len(doc_fitz):
                    img_list = doc_fitz[i].get_images(full=True)
                    for img_idx, img in enumerate(img_list):
                        xref = img[0]
                        base_image = doc_fitz.extract_image(xref)
                        
                        # 提交任务，传入 page_text
                        future = executor.submit(
                            self._describe_image, 
                            base_image["image"], 
                            f"P{i}", 
                            page_text 
                        )
                        image_futures.append(future)
                
                page_tasks.append((page_text, image_futures))

            if doc_fitz: doc_fitz.close()

            # 收集结果
            print(f"  > 等待 {len(page_tasks)} 页的图片解析结果...")
            for i, (text, futures) in enumerate(page_tasks):
                descriptions = ""
                for future in futures:
                    try: descriptions += future.result()
                    except: pass
                
                full_content = f"--- 第 {i+1} 页 ---\n{text}\n{descriptions}"
                results.append({"text": full_content})
        return results

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，支持多模态并发加速"""
        print(f"正在加载 PPT (并发模式): {file_path}")
        results = []
        prs = Presentation(file_path)
        slide_tasks = []

        with concurrent.futures.ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            for i, slide in enumerate(prs.slides):
                slide_text_parts = []
                temp_images = []
                
                # 提取文本和收集图片对象
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        slide_text_parts.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([c.text_frame.text for c in row.cells])
                            slide_text_parts.append(row_text)
                    if shape.shape_type == 13: 
                        temp_images.append(shape)
                
                full_slide_text = "\n".join(slide_text_parts)
                
                # 提交图片任务
                image_futures = []
                for img_shape in temp_images:
                    try:
                        future = executor.submit(
                            self._describe_image, 
                            img_shape.image.blob, 
                            f"S{i}", 
                            full_slide_text
                        )
                        image_futures.append(future)
                    except: pass
                
                slide_tasks.append((full_slide_text, image_futures))

            # 收集结果
            print(f"  > 等待 PPT 图片解析结果...")
            for i, (text, futures) in enumerate(slide_tasks):
                descriptions = ""
                for future in futures:
                    try: descriptions += future.result()
                    except: pass

                final_text = f"--- 幻灯片 {i+1} ---\n{text}\n{descriptions}"
                results.append({"text": final_text})

        return results

    def load_docx(self, file_path: str) -> str:
        print(f"正在加载 DOCX: {file_path}")
        return docx2txt.process(file_path)

    def load_code_or_txt(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        try:
            with open(file_path, 'r', encoding='utf-8') as f: content = f.read()
        except:
            try: 
                with open(file_path, 'r', encoding='gbk') as f: content = f.read()
            except: 
                return ""
            
        # 加上 markdown 标记
        if ext in self.code_map:
            return f"```{self.code_map[ext]}\n{content}\n```"
        return content

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        """
        加载单个文档 (通用入口)
        """
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []
        
        # 1. 根据后缀分发
        if ext == ".pdf": 
            pages = self.load_pdf(file_path)
        elif ext == ".pptx": 
            pages = self.load_pptx(file_path)
        elif ext in [".png", ".jpg", ".jpeg"]: 
            # [新增] 图片处理逻辑
            content = self.load_image(file_path)
            return [{"content": content, "filename": filename, "filetype": ext, "page_number": 1, "summary": "图片文件"}]
        elif ext == ".docx": 
            content = self.load_docx(file_path)
            pages = [{"text": content}] if content else []
        elif ext in self.supported_formats: # 处理所有代码和txt
            content = self.load_code_or_txt(file_path)
            pages = [{"text": content}] if content else []
        else: return []

        # 2. 生成摘要并组装 (针对 PDF, PPTX, DOCX, TXT)
        if ext in [".pdf", ".pptx"]:
            full_txt = "".join([p["text"] for p in pages])
            summary = self._generate_summary(full_txt)
            for i, p in enumerate(pages, 1):
                documents.append({
                    "content": p["text"],
                    "filename": filename,
                    "filepath": file_path,
                    "filetype": ext,
                    "page_number": i,
                    "summary": summary
                })
        elif pages:
            summary = self._generate_summary(pages[0]["text"])
            documents.append({
                "content": pages[0]["text"],
                "filename": filename,
                "filepath": file_path,
                "filetype": ext,
                "page_number": 1,
                "summary": summary
            })
            
        return documents

    def parse_file_to_text(self, file_path: str) -> str:
        """
        [新增] 供前端调用：直接将文件解析为纯字符串，不存库。
        这会触发 load_document，从而触发多模态解析。
        """
        docs = self.load_document(file_path)
        if not docs: return ""
        return "\n\n".join([d["content"] for d in docs])

    def load_all_documents(self) -> List[Dict[str, str]]:
        """批量加载"""
        if not os.path.exists(self.data_dir):
            os.makedirs(self.data_dir, exist_ok=True)
            return []

        documents = []
        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)
        return documents