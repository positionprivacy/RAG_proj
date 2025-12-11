import os
import io
from typing import List, Dict, Optional
import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 引入多模态需要的库
import fitz  # PyMuPDF, 用于PDF图片提取
import dashscope
from http import HTTPStatus
import config  # 导入配置文件

class DocumentLoader:
    def __init__(
        self,
        data_dir: str = config.DATA_DIR,
    ):
        self.data_dir = data_dir
        # 在这里增加了 .py 支持，因为你的作业里包含代码文件
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt", ".py"] 
        
        # 配置 Dashscope API (用于多模态)
        dashscope.api_key = config.OPENAI_API_KEY
        
    def _describe_image(self, image_bytes: bytes, source_info: str, context_text: str = "") -> str:
        """
        [内部辅助函数] 调用 Qwen-VL 对图片进行描述
        增加了 context_text 参数，用于传入当前页面的文字信息
        """
        # ... (前面的代码保持不变: 检查配置, 保存临时文件, 检查文件大小) ...
        if not hasattr(config, "VL_MODEL_NAME") or not config.VL_MODEL_NAME:
            return ""
            
        try:
            # ... (保存临时文件 temp_img_path 的代码保持不变) ...
            temp_img_path = "temp_image_processing.png"
            with open(temp_img_path, "wb") as f:
                f.write(image_bytes)
            
            if os.path.getsize(temp_img_path) < 5 * 1024:
                return ""

            print(f"  > 正在调用 Qwen-VL 理解图片 ({source_info})...")
            
            # --- 修改重点：构建包含上下文的 Prompt ---
            # 截取一部分上下文，防止token溢出（比如取前1000字），通常一页PPT/PDF字数不会太多
            safe_context = context_text[:1000] if context_text else "无"
            
            prompt_content = (
                f"这张图出现在课程课件中。\n"
                f"【该页面的文字内容参考】：{safe_context}\n\n"
                f"请结合上下文详细描述图片内容，提取其中的文字、公式或图表含义。"
                f"如果图片与上下文强相关，请指出它们的关系。"
            )
            # -------------------------------------

            messages = [
                {
                    "role": "user",
                    "content": [
                        {"image": f"file://{os.path.abspath(temp_img_path)}"},
                        {"text": prompt_content} # 这里使用新的 prompt
                    ]
                }
            ]
            
            response = dashscope.MultiModalConversation.call(
                model=config.VL_MODEL_NAME,
                messages=messages
            )
            
            # ... (后续处理 response 的代码保持不变) ...
            if response.status_code == HTTPStatus.OK:
                desc = response.output.choices[0].message.content[0]['text']
                # ...
                return f"\n[图片内容描述]: {desc}\n"
            # ...
            
        except Exception as e:
            # ...
            return ""
        finally:
            if os.path.exists("temp_image_processing.png"):
                os.remove("temp_image_processing.png")

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，支持多模态"""
        print(f"正在加载 PDF: {file_path}")
        results = []
        
        # 1. 使用 PyPDF2 提取文本 (老师要求)
        reader = PdfReader(file_path)
        
        # 2. 使用 PyMuPDF (fitz) 提取图片 (加分项)
        # 为了兼容性，如果没装 pymupdf，这部分会跳过
        doc_fitz = None
        try:
            doc_fitz = fitz.open(file_path)
        except Exception:
            print("提示: 未安装 PyMuPDF，跳过图片识别")

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or "" # 这里已经获取了文本
            
            # --- 多模态图片处理逻辑 Start ---
            image_descriptions = ""
            if doc_fitz and i < len(doc_fitz):
                img_list = doc_fitz[i].get_images(full=True)
                for img_idx, img in enumerate(img_list):
                    xref = img[0]
                    base_image = doc_fitz.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # --- 修改重点：传入 page_text ---
                    desc = self._describe_image(
                        image_bytes, 
                        source_info=f"Page {i+1} Img {img_idx+1}",
                        context_text=page_text  # 传入当前页文本
                    )
                    image_descriptions += desc
            # --- 多模态图片处理逻辑 End ---

            full_content = f"--- 第 {i+1} 页 ---\n{page_text}\n{image_descriptions}"
            results.append({"text": full_content})

        if doc_fitz:
            doc_fitz.close()
            
        return results

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，支持多模态"""
        print(f"正在加载 PPT: {file_path}")
        results = []
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):
            slide_text_parts = []
            images_to_process = [] # 临时存储需要处理的图片对象
            
            # --- 第一轮循环：提取文本 并 收集图片对象 ---
            for shape in slide.shapes:
                # 1. 提取文本框
                if hasattr(shape, "text_frame") and shape.text_frame:
                    slide_text_parts.append(shape.text_frame.text)
                
                # 2. 提取表格
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text_frame.text for cell in row.cells])
                        slide_text_parts.append(row_text)
                
                # 3. 标记图片（先不处理，存起来）
                if shape.shape_type == 13: 
                    images_to_process.append(shape)

            # 组合当前页所有文本
            full_slide_text = "\n".join(slide_text_parts)

            # --- 第二轮循环：集中处理图片（此时已有完整文本） ---
            image_descriptions = []
            for img_shape in images_to_process:
                try:
                    image_bytes = img_shape.image.blob
                    # --- 修改重点：传入 full_slide_text ---
                    desc = self._describe_image(
                        image_bytes, 
                        source_info=f"Slide {i+1}",
                        context_text=full_slide_text
                    )
                    image_descriptions.append(desc)
                except Exception as e:
                    pass 

            # 拼接最终结果：文本 + 图片描述
            final_text = f"--- 幻灯片 {i+1} ---\n{full_slide_text}\n" + "".join(image_descriptions)
            results.append({"text": final_text})

        return results

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件"""
        print(f"正在加载 DOCX: {file_path}")
        # docx2txt 可以处理基本的图片占位，但主要还是提纯文本
        text = docx2txt.process(file_path)
        return text

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件 (同时也用来处理 .py 代码文件)"""
        print(f"正在加载 TXT/Code: {file_path}")
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 如果是 python 代码，给他加上 markdown 标记，方便 LLM 识别
            if file_path.endswith('.py'):
                return f"```python\n{content}\n```"
            
            return content
        except UnicodeDecodeError:
            # 备用编码尝试
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        """
        加载单个文档，保持老师原有逻辑，增加了 .py 的处理
        """
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        # 统一处理逻辑
        content_chunks = []
        
        if ext == ".pdf":
            # load_pdf 返回的是 list[dict]
            pages = self.load_pdf(file_path)
            for i, p in enumerate(pages, 1):
                documents.append({
                    "content": p["text"],
                    "filename": filename,
                    "filepath": file_path,
                    "filetype": ext,
                    "page_number": i
                })
            return documents
            
        elif ext == ".pptx":
            # load_pptx 返回的是 list[dict]
            slides = self.load_pptx(file_path)
            for i, s in enumerate(slides, 1):
                documents.append({
                    "content": s["text"],
                    "filename": filename,
                    "filepath": file_path,
                    "filetype": ext,
                    "page_number": i
                })
            return documents
        
        elif ext == ".docx":
            content = self.load_docx(file_path)
        elif ext == ".txt" or ext == ".py":
            content = self.load_txt(file_path)
        else:
            print(f"不支持的文件格式: {ext}")
            return []

        # 对于 docx, txt, py，作为一个整体块返回（或者你可以在这里做初步分割）
        if content:
            documents.append({
                "content": content,
                "filename": filename,
                "filepath": file_path,
                "filetype": ext,
                "page_number": 1 
            })
            
        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        """加载数据目录下的所有文档"""
        if not os.path.exists(self.data_dir):
            os.makedirs(self.data_dir, exist_ok=True) # 防止目录不存在报错
            print(f"数据目录不存在，已创建空目录: {self.data_dir}")
            return []

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        print(f"所有文档加载完毕，共处理 {len(documents)} 个片段。")
        return documents